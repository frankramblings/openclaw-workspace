"""OpenClaw Workspace — FastAPI app.

Serves the (reused) Odysseus SPA and wires:
  - /api/chat_stream  → the bridge to OpenClaw's gateway brain  (REAL, v1)
  - /api/items        → native unified inbox (gmail/slack/asana/obsidian collectors)
  - a handful of minimal stubs so the SPA loads without console errors

Run:  uvicorn backend.app:app --reload --port 8800   (from the repo root)
"""
from __future__ import annotations

import asyncio
import base64
import contextlib
import json
import logging
import mimetypes
import re
import subprocess
import tempfile
import time
from contextlib import asynccontextmanager
from pathlib import Path

from fastapi import Body, FastAPI, Form
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from starlette.middleware.gzip import GZipMiddleware

from . import branch_context, bridge, capabilities, chat_search, config, doctor, draft_mode, event_store, followup, monitor, sessions_store, terminals, websearch
from .auth_gate import AuthGateMiddleware
from .memory import maybe_auto_extract
from .calendar import router as calendar_router
from .cron import router as cron_router
from .documents import router as documents_router
from .email_himalaya import router as email_router
from .emoji_proxy import router as emoji_router
from .inbox import router as inbox_router
from .jobs import router as jobs_router
from .memory import router as memory_router
from .notes import router as notes_router
from .research import router as research_router
from .settings_status import router as settings_router
from .skills import router as skills_router
from .uploads import ATTACH_DIR
from .uploads import router as uploads_router
from .workspace_files import router as workspace_files_router
from .workspace_watch import router as workspace_watch_router
from . import workspace_watch
from .terminals import router as terminals_router
from .resume_route import router as resume_router
from .export_pdf import router as export_pdf_router
from . import workspace_files

async def _startup_reindex() -> None:
    """Keep the semantic-search index fresh in the background. Runs the first
    build shortly after boot (delayed so the app serves requests immediately),
    then refreshes every 30 min. reindex() is incremental — sessions whose
    `updated` stamp is unchanged are skipped — so each refresh only embeds new
    or changed conversations. Failures are swallowed (log only) so a
    gateway/Voyage hiccup can never crash boot or stop the loop."""
    log = logging.getLogger("workspace.chat_search")
    await asyncio.sleep(20)
    while True:
        try:
            await chat_search.reindex()
        except asyncio.CancelledError:
            raise
        except Exception:  # noqa: BLE001 - never let an index failure stop the loop
            log.warning("periodic reindex failed", exc_info=True)
        await asyncio.sleep(1800)  # 30 min — pick up new/updated conversations


@asynccontextmanager
async def _lifespan(_app: FastAPI):
    # The persistent gateway monitor (status dot / restart awareness).
    task = asyncio.create_task(monitor.run())
    # Non-blocking semantic-search index build (delayed; swallows failures).
    search_task = asyncio.create_task(_startup_reindex())
    # Followup promises: deadline + crash-recovery backstop.
    followup_task = asyncio.create_task(followup.sweeper())
    # Filesystem watcher for the doc editor's live-refresh (broadcasts to
    # /api/workspace/watch subscribers). Cheap Rust-backed inotify; one task.
    workspace_watch.start_watcher()
    yield
    # Reap every PTY shell so its `bash -i` child (and descendants) get
    # SIGHUP→SIGKILL right now. Otherwise they ignore the SIGTERM systemd sends
    # the whole cgroup, holding the unit open until TimeoutStopSec forces a
    # SIGKILL of everything — the intermittent restart hang / 502 window. We call
    # PtySession.close() directly (not close_session) so we DON'T also unlink the
    # persistent per-session image-attachment registry; scrollback/cwd are
    # already flushed by close() and restored on the next attach.
    for _sess in list(getattr(terminals, "_sessions", {}).values()):
        try:
            _sess.close()
        except Exception:  # noqa: BLE001 - best-effort teardown, never block exit
            pass
    task.cancel()
    search_task.cancel()
    followup_task.cancel()
    for t in (task, search_task, followup_task):
        with contextlib.suppress(asyncio.CancelledError):
            await t


app = FastAPI(title="OpenClaw Workspace", lifespan=_lifespan)

# Wire bytes matter on the phone-over-Tailscale path and nothing upstream
# compresses (Tailscale Serve passes bytes through): style.css alone is 1MB
# raw / 227KB gzipped. Streaming responses (SSE) are flushed per-chunk by
# Starlette's GZipResponder, so /api/chat/stream keeps streaming.
app.add_middleware(GZipMiddleware, minimum_size=1024)

# Optional auth gate: a complete no-op when WORKSPACE_AUTH_TOKEN is unset
# (the middleware reads the token at request time and short-circuits immediately
# when it's None, so zero overhead or behavior change for the default no-token
# case). Added AFTER GZip so auth runs in the outer layer (before compression).
app.add_middleware(AuthGateMiddleware)

app.include_router(inbox_router)
app.include_router(jobs_router)
app.include_router(memory_router)
app.include_router(skills_router)
app.include_router(cron_router)
app.include_router(email_router)
app.include_router(calendar_router)
app.include_router(settings_router)
app.include_router(notes_router)
app.include_router(documents_router)
app.include_router(followup.router)
app.include_router(uploads_router)
app.include_router(research_router)
app.include_router(emoji_router)
app.include_router(workspace_files_router)
app.include_router(workspace_watch_router)
app.include_router(terminals_router)
app.include_router(resume_router)
app.include_router(export_pdf_router)

# Active gateway runs by sessionKey, so the Stop button can chat.abort the run
# server-side. chat.js already POSTs /api/chat/stop/<sid> on explicit Stop
# (abortCurrentRequest(true)) — until now that hit the GET-only catch-all and
# only the browser-side fetch died, while the codex run kept burning.
_ACTIVE_RUNS: dict[str, dict] = {}

# Fire-and-forget background tasks (memory auto-extract, gateway-side session
# delete, on-demand reindex). asyncio holds only a WEAK reference to a bare
# create_task(), so a long one can be garbage-collected mid-flight; keep a
# strong ref here and drop it when it finishes.
_BG_TASKS: set[asyncio.Task] = set()


def _spawn(coro) -> asyncio.Task:
    """create_task + keep a strong reference until the task completes."""
    task = asyncio.create_task(coro)
    _BG_TASKS.add(task)
    task.add_done_callback(_BG_TASKS.discard)
    return task


# session_key -> the detached asyncio.Task that drives ONE turn's frames into
# the resumable event_store. The recorder is deliberately decoupled from any
# browser: the gateway agent keeps working (and we keep recording) even after
# the reader that started the turn refreshes, switches threads, or closes the
# tab. Every reader — the original POST, a thread-switch return, a post-reload
# resume — is just a tail of event_store, so a dropped reader can never stop a
# turn or lose the work done while away.
_TURN_TASKS: dict[str, asyncio.Task] = {}

# Idle keepalive for the POST tail (seconds): a `: keepalive` comment keeps
# proxies from killing the connection while the agent is thinking. Matches
# resume_route's tail loop.
_TURN_KEEPALIVE_S = 15.0


async def _record_turn(session_key: str, source) -> None:
    """The single writer. Drain `source` (the async generator of SSE frames for
    one turn) into `event_store`, independent of any reader. Marks the turn
    active on entry and inactive on exit, and ALWAYS lands a terminal `[DONE]`
    frame so every tail closes cleanly — even on error or an explicit Stop
    (CancelledError). Appending must never raise into the turn."""
    event_store.begin_turn(session_key)
    done_emitted = False
    try:
        async for chunk in source:
            try:
                event_store.append(session_key, chunk)
            except Exception:  # noqa: BLE001 - event log issue can't break the turn
                pass
            if _is_done_frame(chunk):
                done_emitted = True
    finally:
        if not done_emitted:
            try:
                event_store.append(session_key, _DONE_SSE)
            except Exception:  # noqa: BLE001
                pass
        event_store.end_turn(session_key)
        _TURN_TASKS.pop(session_key, None)


def _start_turn_recorder(session_key: str, source_factory):
    """Launch the detached recorder for a turn if one isn't already running for
    this session. `source_factory` is a zero-arg callable returning the turn's
    SSE async generator; it is invoked ONLY when we actually start, so a guarded
    no-op never spins up a second gateway run. Returns the recorder Task."""
    prev = _TURN_TASKS.get(session_key)
    if prev is not None and not prev.done():
        return prev  # a turn is already recording for this session
    task = asyncio.create_task(_record_turn(session_key, source_factory()))
    _TURN_TASKS[session_key] = task
    return task


@app.get("/api/health")
async def health():
    return {
        "ok": True,
        "gateway": config.gateway_ws_url(),
        "session": config.session_key(),
        "has_password": bool(config.gateway_password()),
    }


@app.get("/api/doctor")
async def api_doctor():
    """Diagnose the OpenClaw connection (read-only)."""
    return doctor.summarize(await doctor.run_checks())


@app.get("/api/config")
async def workspace_config():
    """Public branding/config the SPA reads at boot so the agent name is right
    even before a frontend re-sync. Single source of truth lives in config.py."""
    return {
        "agent_name": config.agent_name(),
        "accent": config.accent_color(),
        # The footer shows this; previously it fetched the ENTIRE workspace
        # tree walk just to read .root (2026-06-12 mobile review E2).
        "workspace_root": str(workspace_files.workspace_root()),
        # AGPL-3.0 §13: the SPA renders a "Source" link to this URL.
        "source_url": config.source_url(),
    }


@app.get("/api/capabilities")
async def api_capabilities():
    """Which tabs are usable on this install (drives UI gating)."""
    return capabilities.snapshot()


@app.get("/api/gateway/status")
async def gateway_status():
    """Last-known gateway state from the persistent monitor, for the UI's
    status dot (polled ~30s). state: ok | restarting | down."""
    return await monitor.status()


# --- The one real, load-bearing endpoint: chat ------------------------------

def _model_ref(rec: dict | None) -> str | None:
    """Build a gateway model ref ("provider/model") from a session record, or
    None to leave the model at the agent default. Returns None for the
    "openclaw" placeholder (legacy/bootstrap) AND when the pick already equals
    the configured default — so we only set an override when it actually
    differs (no per-turn sessions.create churn for default chats)."""
    if not rec:
        return None
    model = (rec.get("model") or "").strip()
    if not model or model == "openclaw":
        return None
    provider = (rec.get("endpoint_id") or "").strip()
    def_provider, def_model = config.default_model()
    if model == def_model and (not provider or provider in (def_provider, "openclaw")):
        return None
    return f"{provider}/{model}" if provider and provider != "openclaw" else model


# --- Auto-title: name a fresh thread from its first message (AI + fallback) ---
# The SPA names new chats "{model} {time}" (a placeholder). On the first message
# we retitle the session — instantly to a first-line snippet, then upgraded to an
# AI title generated by the brain (ChatGPT-style). Backend-only: the Library list
# already reloads after a turn, so the new title just appears.

_SPEED_THINKING = {"fast": "low", "deep": "high"}


def _thinking_for_speed(speed: str | None) -> str | None:
    """Map the chat's speed setting to chat.send's per-turn thinking override.
    normal (and anything unknown) sends NO override — the default path stays
    byte-identical to pre-toggle behavior."""
    return _SPEED_THINKING.get(speed or "")


_DONE_SSE = "data: [DONE]\n\n"


def _is_done_frame(chunk: str) -> bool:
    """True only for the exact terminal marker `data: [DONE]` — NOT for a delta
    whose text merely CONTAINS the literal "[DONE]". Frames are single
    `data: <body>` SSE messages; comparing the stripped body exactly stops a
    real reply (or a message about this very code) that mentions [DONE] from
    being dropped mid-stream or cutting the tail short."""
    for line in chunk.splitlines():
        if line.startswith("data:") and line[5:].strip() == "[DONE]":
            return True
    return False


_TITLE_SESSION_KEY = f"{config.web_session_prefix()}-titler"
# "{base} 1:56:53 PM" / "{base} 14:05:09" — the SPA's placeholder name.
_PLACEHOLDER_RE = re.compile(r".+\s\d{1,2}:\d{2}:\d{2}(\s?[AP]M)?$", re.I)


def _needs_title(rec: dict) -> bool:
    name = (rec.get("name") or "").strip()
    if not name or name in ("New chat", "Nobody"):
        return True
    return bool(_PLACEHOLDER_RE.match(name))


def _first_chars_title(message: str, n: int = 42) -> str:
    text = (message or "").strip()
    if not text:
        return ""
    line = text.splitlines()[0].strip()
    return line if len(line) <= n else line[:n].rstrip() + "…"


def _sanitize_title(raw: str) -> str:
    if not raw:
        return ""
    line = raw.strip().splitlines()[0].strip().strip('"\'').strip()
    line = re.sub(r"^(chat title|title)\s*[:\-]\s*", "", line, flags=re.I)
    return line.rstrip(" .!,;:")[:60].strip()


async def _collect_brain_text(prompt: str, session_key: str,
                              model_ref: str | None = None) -> str:
    chunks: list[str] = []
    async for sse in bridge.stream_turn(prompt, session_key=session_key,
                                        model_ref=model_ref):
        if not sse.startswith("data:"):
            continue
        body = sse[5:].strip()
        if not body or body == "[DONE]":
            continue
        try:
            obj = json.loads(body)
        except Exception:  # noqa: BLE001
            continue
        if isinstance(obj, dict) and obj.get("delta"):
            chunks.append(obj["delta"])
    return "".join(chunks).strip()


async def _generate_ai_title(message: str) -> str:
    prompt = ("Generate a short, specific chat title (3-6 words, no quotes, no "
              "trailing punctuation) for the topic of a conversation that opens "
              f"with this message:\n\n{message[:600]}\n\nOutput ONLY the title.")
    return _sanitize_title(await _collect_brain_text(
        prompt, _TITLE_SESSION_KEY, model_ref=config.TITLE_MODEL))


def _wants_web(use_web: str, allow_web_search: str) -> bool:
    """The globe toggle's FIELD NAME depends on the SPA's vestigial chat/agent
    mode: `use_web` in chat mode, `allow_web_search` in agent mode (chat.js
    ~747-753). Accept both — agent mode silently lost web search before."""
    return any(v.lower() in ("true", "1", "on", "yes")
               for v in (use_web, allow_web_search))


def _sse_frame(chunk: str):
    """Parse one of the bridge's `data: {...}` SSE strings back to its dict
    (None for [DONE]/non-JSON). Used to observe what the turn relayed."""
    body = chunk[5:].strip() if chunk.startswith("data:") else ""
    if not body or body == "[DONE]":
        return None
    try:
        obj = json.loads(body)
    except ValueError:
        return None
    return obj if isinstance(obj, dict) else None


def reply_after(history: list, brain_message: str) -> str | None:
    """Assistant text following the LAST occurrence of this turn's user message
    in the transcript — i.e. THIS turn's reply, never an earlier one."""
    idx = next((i for i in range(len(history) - 1, -1, -1)
                if history[i].get("role") == "user"
                and history[i].get("content") == brain_message), None)
    if idx is None:
        return None
    text = "\n".join(str(m.get("content") or "") for m in history[idx + 1:]
                     if m.get("role") == "assistant").strip()
    return text or None


def _turn_timing_record(run_info: dict, session_key: str, model_ref: str | None,
                        *, text_seen: bool, failed: bool,
                        thinking: str | None = None) -> dict:
    """One flat JSONL record describing where this turn's wall-clock went.
    All *_ms fields are measured from chat.send write; None = never happened."""
    timing = run_info.get("timing") or {}

    def ms(a: str, b: str) -> int | None:
        return (int((timing[b] - timing[a]) * 1000)
                if a in timing and b in timing else None)

    return {
        "ts": int(time.time()),
        "session": session_key,
        "model": model_ref or "default",
        "thinking": thinking,   # the chat.send override sent (None = normal)
        "ack_ms": ms("t_send", "t_ack"),
        "first_frame_ms": ms("t_send", "t_first_frame"),
        "first_text_ms": ms("t_send", "t_first_text"),
        "late_ms": ms("t_send", "t_late"),
        "total_ms": ms("t_send", "t_end"),
        "stalled": bool(run_info.get("stalled")),
        "retried": bool(run_info.get("retried")),
        "text_seen": text_seen,
        "failed": failed,
    }


def _log_turn_timing(record: dict) -> None:
    """Append one JSONL line to .data/turn_timings.jsonl. Telemetry must never
    break a turn: every failure is swallowed. Single-generation rotation at
    2MB keeps the file bounded on this disk-starved box."""
    with contextlib.suppress(Exception):
        path = config.DATA_DIR / "turn_timings.jsonl"
        path.parent.mkdir(parents=True, exist_ok=True)
        if path.exists() and path.stat().st_size > 2_000_000:
            path.replace(path.with_name(path.name + ".old"))
        with path.open("a", encoding="utf-8") as fh:
            fh.write(json.dumps(record) + "\n")


# First checks fire fast — the reply usually already sits in the transcript
# when this poll starts (it lands seconds *before* we get here on slow turns,
# milliseconds after on fast ones). Tail stays ~10s total like the old
# 5 × 2s schedule.
_LATE_REPLY_SCHEDULE = (0.3, 0.5, 1.0, 2.0, 2.0, 2.0, 2.2)


async def _late_reply(session_key: str, brain_message: str,
                      _sleep=asyncio.sleep) -> str | None:
    """Fetch the reply that the gateway commits to the transcript only AFTER
    the run's lifecycle end (message-tool delivery — see _relay_events docs).
    Polls with fast-start backoff; returns None if nothing lands (genuinely
    textless turn)."""
    for delay_s in _LATE_REPLY_SCHEDULE:
        await _sleep(delay_s)
        try:
            data = await bridge.fetch_history(session_key)
        except Exception:  # noqa: BLE001 - transient WS trouble: keep polling
            continue
        text = reply_after(data.get("history") or [], brain_message)
        if text:
            return text
    return None


_log = logging.getLogger(__name__)

# Gateway rejects chat.send payloads over ~10 MB base64. Cap individual images
# at 4 MB raw (~5.3 MB base64) to stay safely under that limit with multiple
# attachments. HEIC/HEIF files are converted to JPEG first (models don't support
# them natively and iPhone photos are often 7+ MB).
_ATTACH_MAX_BYTES = 4 * 1024 * 1024  # 4 MB


def _heic_to_jpeg(src: Path) -> bytes | None:
    """Convert a HEIC/HEIF file to JPEG via ffmpeg. Returns JPEG bytes or None."""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
            tmp_path = tmp.name
        result = subprocess.run(
            ["ffmpeg", "-y", "-i", str(src), "-update", "1",
             "-vf", "scale='min(1920,iw)':-2", "-q:v", "5", tmp_path],
            capture_output=True, timeout=30,
        )
        if result.returncode == 0:
            return Path(tmp_path).read_bytes()
    except Exception:  # noqa: BLE001
        pass
    finally:
        if tmp_path:
            with contextlib.suppress(OSError):
                Path(tmp_path).unlink()
    return None


def _resolve_attachments(raw: str) -> list[dict]:
    """Turn the composer's posted `attachments` (a JSON array of upload ids, each
    a filename under ATTACH_DIR) into the chat.send attachment shape the gateway
    accepts: {type, mimeType, fileName, content(base64)}. Only image/* files are
    forwarded — the gateway sniffs the bytes and drops non-images, and gpt-5.5
    takes them as inline vision blocks (large ones it offloads to media refs the
    agent resolves). Silently skips anything missing/unreadable so a bad id never
    breaks the turn."""
    if not raw:
        return []
    try:
        ids = json.loads(raw)
    except Exception:  # noqa: BLE001 - malformed field → no attachments
        return []
    out: list[dict] = []
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")  # path-traversal guard
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        if not mime.startswith("image/"):
            continue  # gateway accepts only image/* attachments
        try:
            data = path.read_bytes()
        except Exception:  # noqa: BLE001 - unreadable file → skip it
            continue
        # Convert HEIC/HEIF → JPEG (models don't support these formats).
        if mime in ("image/heic", "image/heif"):
            converted = _heic_to_jpeg(path)
            if converted:
                data, mime = converted, "image/jpeg"
            else:
                _log.warning("HEIC conversion failed for %s; skipping attachment", path.name)
                continue
        # Skip images that are still too large for the gateway after conversion.
        if len(data) > _ATTACH_MAX_BYTES:
            _log.warning("Attachment %s is %d bytes (>4 MB); skipping", path.name, len(data))
            continue
        out.append({
            "type": "image",
            "mimeType": mime,
            "fileName": path.stem + (".jpg" if mime == "image/jpeg" else path.suffix),
            "content": base64.b64encode(data).decode("ascii"),
        })
    return out


# Non-image files: extract text and inline into the message body. Total across
# all files per turn is capped so a huge upload can't blow the context.
_TEXT_ATTACH_TOTAL_MAX = 200 * 1024
_TEXT_ATTACH_PER_FILE_MAX = 100 * 1024
_TEXT_RAW_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log",
                  ".xml", ".yaml", ".yml", ".ini", ".conf", ".toml", ".env",
                  ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".htm", ".css",
                  ".scss", ".sh", ".bash", ".zsh", ".fish", ".rb", ".go", ".rs",
                  ".java", ".kt", ".swift", ".c", ".h", ".cpp", ".hpp", ".cc",
                  ".cs", ".php", ".sql", ".rst", ".org", ".vue", ".svelte",
                  ".lua", ".pl", ".r", ".dockerfile", ".makefile", ".gradle",
                  ".diff", ".patch"}


def _extract_file_text(path: Path, mime: str) -> str | None:
    """Best-effort text extraction. Returns None for formats we can't read (e.g.
    .numbers without libreoffice); the caller silently skips those."""
    ext = path.suffix.lower()
    try:
        if ext in _TEXT_RAW_EXTS or mime.startswith("text/") or mime == "application/json":
            return path.read_text(errors="replace")
        if ext == ".xlsx" or mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ):
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
            wb.close()
            return "\n".join(parts)
        if ext == ".pdf" or mime == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            out: list[str] = []
            for i, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    out.append(f"## Page {i}\n{text}")
            return "\n\n".join(out)
        if ext == ".docx" or mime == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            from docx import Document
            doc = Document(str(path))
            parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
            for tbl in doc.tables:
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
            return "\n".join(parts)
        if ext == ".pptx" or mime == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            from pptx import Presentation
            pres = Presentation(str(path))
            parts: list[str] = []
            for i, slide in enumerate(pres.slides, start=1):
                parts.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(r.text for r in para.runs).strip()
                            if line:
                                parts.append(line)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            if any(cells):
                                parts.append("\t".join(cells))
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Speaker notes] {notes}")
            return "\n".join(parts)
    except Exception as e:  # noqa: BLE001
        _log.warning("Text extraction failed for %s (%s): %s", path.name, ext, e)
        return None
    return None


def _extract_text_attachments(raw: str) -> tuple[list[dict], list[str]]:
    """Return (files, skipped). `files` is [{name, text}] for uploads we could
    read; `skipped` is names of non-image uploads whose content we couldn't
    extract (unsupported format, empty output, or budget exceeded). Images are
    handled by `_resolve_attachments` and are neither returned nor skipped
    here."""
    if not raw:
        return [], []
    try:
        ids = json.loads(raw)
    except Exception:  # noqa: BLE001
        return [], []
    out: list[dict] = []
    skipped: list[str] = []
    total = 0
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or ""
        if mime.startswith("image/"):
            continue  # image path handles these
        text = _extract_file_text(path, mime)
        if not text or not text.strip():
            skipped.append(path.name)
            continue
        text = text.strip()
        if len(text) > _TEXT_ATTACH_PER_FILE_MAX:
            text = text[:_TEXT_ATTACH_PER_FILE_MAX] + "\n... [truncated]"
        if total + len(text) > _TEXT_ATTACH_TOTAL_MAX:
            _log.warning("Text-attach total exceeded %d bytes; dropping %s",
                         _TEXT_ATTACH_TOTAL_MAX, path.name)
            skipped.append(path.name)
            continue
        total += len(text)
        out.append({"name": path.name, "text": text})
    return out, skipped


def _prepend_text_attachments(message: str, files: list[dict],
                              skipped: list[str] | None = None) -> str:
    """Inline extracted file text into the user message so it reaches the model
    regardless of vision/file-type support. The gateway only forwards image
    attachments; everything else has to arrive as text."""
    skipped = skipped or []
    if not files and not skipped:
        return message
    parts: list[str] = []
    if files:
        parts.append("The user attached the following file(s):\n")
        for f in files:
            parts.append(f"── {f['name']} ──")
            parts.append(f["text"])
            parts.append("")
        parts.append("──\n")
    if skipped:
        parts.append(
            "[System note: couldn't extract text from these attachments — "
            "unsupported format or empty content: "
            + ", ".join(skipped) + ". Tell the user which files were dropped "
            "so they can convert or resend.]\n"
        )
    parts.append(message or "")
    return "\n".join(parts)


_CHAT_ATTACH_DIR = ATTACH_DIR.parent / ".chat-attachments"


def _attach_log_path(session_id: str) -> Path | None:
    safe = "".join(c for c in (session_id or "") if c.isalnum() or c in "-_.")
    return (_CHAT_ATTACH_DIR / f"{safe}.json") if safe else None


def _persist_msg_attachments(session_id: str, message: str, attachments_raw: str) -> None:
    """Record image attachments for a sent message so /api/history can rehydrate
    them on reload — the gateway transcript keeps only text, so a sent image
    otherwise vanishes on refresh. Appends one record per send-with-images:
    {text, att:[{id,url}]}. Best-effort; never raises into the turn."""
    if not session_id or not attachments_raw:
        return
    try:
        ids = json.loads(attachments_raw)
    except Exception:  # noqa: BLE001 - malformed field → nothing to persist
        return
    att: list[dict] = []
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or ""
        if not mime.startswith("image/"):
            continue
        att.append({"id": safe, "url": f"/api/upload/{safe}"})
    if not att:
        return
    p = _attach_log_path(session_id)
    if not p:
        return
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        log = []
        if p.exists():
            try:
                log = json.loads(p.read_text() or "[]")
            except Exception:  # noqa: BLE001
                log = []
        log.append({"text": (message or "").strip(), "att": att})
        p.write_text(json.dumps(log))
    except Exception:  # noqa: BLE001 - persistence is best-effort
        pass


def _apply_msg_attachments(session_id: str, history: list[dict]) -> None:
    """Assign persisted image attachments to user messages in `history`, matching
    on (trimmed) text in send order. Mutates `history` in place; best-effort."""
    p = _attach_log_path(session_id)
    if not p or not p.exists():
        return
    try:
        log = json.loads(p.read_text() or "[]")
    except Exception:  # noqa: BLE001
        return
    if not isinstance(log, list) or not log:
        return
    used = [False] * len(log)
    for m in history:
        if m.get("role") != "user":
            continue
        text = (m.get("content") or "").strip()
        for i, rec in enumerate(log):
            if used[i]:
                continue
            if (rec.get("text") or "").strip() == text:
                m["attachments"] = rec.get("att") or []
                used[i] = True
                break


def _terminal_attachments(terminal_key: str) -> list[dict]:
    """Pending images the user dropped into this chat's terminal → chat.send
    image blocks, then mark them consumed so each rides exactly one turn. The
    token→path mapping itself persists (terminals registry) for later resolves.
    Mirrors _resolve_attachments' block shape; image/* only; bad files skipped."""
    out: list[dict] = []
    consumed: list[str] = []
    for it in terminals.list_attachments(terminal_key, pending_only=True):
        path = Path(it.get("path", ""))
        if not path.is_file():
            continue
        mime = it.get("mime") or mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        if not mime.startswith("image/"):
            continue
        try:
            data = path.read_bytes()
        except Exception:  # noqa: BLE001 - unreadable → skip, never break the turn
            continue
        out.append({
            "type": "image",
            "mimeType": mime,
            "fileName": path.name,
            "content": base64.b64encode(data).decode("ascii"),
        })
        consumed.append(it["token"])
    if consumed:
        terminals.mark_consumed(terminal_key, consumed)
    return out


@app.post("/api/chat_stream")
async def chat_stream(message: str = Form(...), session: str = Form(default=""),
                      use_web: str = Form(default=""),
                      allow_web_search: str = Form(default=""),
                      attachments: str = Form(default=""),
                      active_doc_id: str = Form(default="")):
    """Stream a turn from OpenClaw's brain as Odysseus-shaped SSE.

    The posted `session` is the SPA's session id; we resolve it to that chat's
    own gateway sessionKey (agent:main:web-<id>) so each Library chat is an
    isolated thread and none contend with Signal on agent:main:main. Unknown
    ids fall back to the shared web key. The session's picked model (if any) is
    applied to that session only, so the picker actually switches the model.

    On a fresh thread's first message we also auto-title it (see above).
    """
    rec = sessions_store.get(session) if session else None
    session_key = rec["sessionKey"] if rec else config.web_session_key()

    # One turn per session at a time. The detached recorder already guards
    # against a second concurrent turn for this sessionKey, so if we proceeded
    # here the new message's _drive_turn would never run — the gateway would
    # never see it and this POST would silently tail the PREVIOUS turn as if it
    # were the answer. Tell the client instead (a fresh chat and a second unsaved
    # chat both share config.web_session_key(), so this also covers double-send).
    prev_turn = _TURN_TASKS.get(session_key)
    if prev_turn is not None and not prev_turn.done():
        async def _busy():
            yield bridge._sse({"type": "tool_start", "tool": "bridge",
                               "tool_id": "busy", "command": "turn in progress",
                               "round": 1})
            yield bridge._sse({"type": "tool_output", "tool": "bridge",
                               "tool_id": "busy", "exit_code": 1,
                               "output": "A turn is already running for this chat — "
                                         "wait for it to finish, then resend."})
            yield _DONE_SSE
        return StreamingResponse(_busy(), media_type="text/event-stream")

    run_info: dict = {}  # bridge fills sessionKey/runId once chat.send acks
    chat_attachments = _resolve_attachments(attachments)  # image uploads → vision
    # Non-image uploads (CSV, XLSX, DOCX, PPTX, PDF, text): extract and inline
    # into the message so they reach the model (the gateway drops non-image
    # bytes). Skipped files are surfaced as a system note so the assistant can
    # tell the user which files were dropped.
    text_files, skipped_files = _extract_text_attachments(attachments)
    if text_files or skipped_files:
        message = _prepend_text_attachments(message, text_files, skipped_files)
    # Persist the image refs (keyed by the SPA session id) so they survive a
    # reload — the gateway transcript only keeps the user's text.
    if session and attachments:
        _persist_msg_attachments(session, message, attachments)

    # Draft mode: chat.js posts active_doc_id whenever the document panel is
    # open (auto-saving the doc first). Snapshot now (the user's undo), wrap
    # the message in gen(), detect agent edits after the turn (draft_mode.py).
    draft_doc = draft_mode.pre_turn(active_doc_id) if active_doc_id else None

    title_task = None
    if rec and message.strip() and _needs_title(rec):
        snippet = _first_chars_title(message)
        if snippet:  # instant fallback so the title is never the timestamp
            sessions_store.update(rec["id"], name=snippet)
        title_task = asyncio.create_task(_generate_ai_title(message))

    async def _drive_turn():
        """Produce every SSE frame for this turn (web search, the gateway relay,
        late reply, metrics, draft doc_update, DONE). This is the turn's source
        of truth; the detached recorder (`_record_turn`) drains it into
        event_store and owns the turn boundary + the event ids. Readers never
        consume this directly — they tail event_store — so the turn survives any
        reader leaving."""
        brain_message = message
        text_seen = False    # any non-thinking {"delta"} relayed this turn?
        tools_seen = False   # any tool card relayed (fresh bubble needed)?
        failed = False       # bridge/agent error or user abort — no reply coming
        try:
            # Web search (the composer's globe toggle — field name varies by
            # the SPA's vestigial mode, see _wants_web). Search failures
            # degrade to a visible failed tool card — never break the turn.
            if _wants_web(use_web, allow_web_search):
                s = websearch.load_settings()
                if s.get("search_provider") != "disabled":
                    yield bridge._sse({"type": "tool_start", "tool": "web_search",
                                       "tool_id": "websearch",
                                       "command": message[:120], "round": 1})
                    try:
                        results = await websearch.search(
                            message, int(s.get("search_result_count") or 5))
                        yield bridge._sse({
                            "type": "tool_output", "tool": "web_search",
                            "tool_id": "websearch", "exit_code": 0,
                            "output": "\n".join(f"[{i+1}] {r['title']} — {r['url']}"
                                                for i, r in enumerate(results))
                                      or "no results"})
                        if results:
                            brain_message = websearch.context_block(message, results)
                    except Exception as exc:  # noqa: BLE001
                        yield bridge._sse({"type": "tool_output", "tool": "web_search",
                                           "tool_id": "websearch",
                                           "output": f"web search failed: {exc}",
                                           "exit_code": 1})
            if draft_doc is not None:
                brain_message = draft_mode.wrap_message(brain_message, draft_doc)
            # Gary-drive: when terminal control is on for this chat, prepend a
            # per-turn capability hint + a freshly-minted token (stripped from
            # the history view in the /api/history handler below).
            # CRITICAL: bind Gary to the SAME PTY key the human panel uses — the
            # SPA chat id (rec["id"], == frontend getCurrentSessionId()), NOT the
            # gateway sessionKey. Otherwise Gary drives a different terminal than
            # the one the user sees. New/unsaved chats fall back to "global",
            # which is what the panel sends (curSession() || "global").
            terminal_key = rec["id"] if rec else "global"
            # Terminal image drops: prepend the (history-stripped) token→path map
            # and merge any pending dropped images into THIS turn's vision blocks.
            att_note = terminals.terminal_attachment_note(terminal_key)
            if att_note:
                brain_message = att_note + brain_message
            turn_attachments = chat_attachments + _terminal_attachments(terminal_key)
            if terminals.gary_mode_for_session(terminal_key):
                brain_message = terminals.gary_capability_note(terminal_key) + brain_message
            _ACTIVE_RUNS[session_key] = run_info

            async for chunk in bridge.stream_turn(brain_message, session_key=session_key,
                                                  model_ref=_model_ref(rec),
                                                  attachments=turn_attachments,
                                                  run_info=run_info,
                                                  thinking=_thinking_for_speed((rec or {}).get("speed"))):
                if _is_done_frame(chunk):
                    continue  # hold DONE until the title settles, then send our own
                frame = _sse_frame(chunk)
                if isinstance(frame, dict):
                    if frame.get("delta") and not frame.get("thinking"):
                        text_seen = True
                    if frame.get("type") in ("tool_start", "tool_output"):
                        tools_seen = True
                    if (frame.get("exit_code") == 1
                            and frame.get("tool") in ("bridge", "agent")
                            # stall terminal card still allows the late-reply
                            # poll to salvage a transcript-landed reply.
                            and frame.get("tool_id") != "stall") \
                            or frame.get("tool_id") == "abort":
                        failed = True
                # Pure source: just yield the frame. The detached recorder
                # (`_record_turn`) is the single writer to event_store; the POST
                # response and every other reader are tails of that log, so a
                # dropped reader can't stop the turn or lose mid-flight frames.
                yield chunk
            # Late delivery: the agent often replies via its `message` tool,
            # whose text lands in the transcript seconds AFTER the run's
            # lifecycle end — the live stream then carries no text at all and
            # the reply only appeared on the next refresh. Same workaround the
            # research engine needed: poll chat.history briefly and emit it.
            if not text_seen and not failed:
                late = await _late_reply(session_key, brain_message)
                if late:
                    run_info.setdefault("timing", {})["t_late"] = time.monotonic()
                    if tools_seen:
                        yield bridge._sse({"type": "agent_step"})  # fresh bubble
                    yield bridge._sse({"delta": late})
            # Final turn metrics: the vendor SPA renders a footer time from a
            # {type:"metrics"} frame (chatRenderer.displayMetrics) that its
            # original backends sent and we never did. response_time prefers
            # t_late over t_end — that's when the user actually saw the reply.
            # No token/cost fields: the gateway doesn't expose usage, and
            # displayMetrics degrades to a plain "12.3s" without them.
            timing = run_info.get("timing") or {}
            if "t_send" in timing and not failed:
                t_done = timing.get("t_late") or timing.get("t_end") or time.monotonic()
                data = {"response_time": round(t_done - timing["t_send"], 1)}
                if "t_first_text" in timing:  # pre-text wait = thinking + prep
                    data["agent_model_wait_time"] = round(
                        timing["t_first_text"] - timing["t_send"], 1)
                if _model_ref(rec):
                    data["model"] = _model_ref(rec)
                yield bridge._sse({"type": "metrics", "data": data})
        finally:
            _ACTIVE_RUNS.pop(session_key, None)
            # begin_turn/end_turn + event_store.append are owned by the detached
            # recorder (_record_turn), NOT this source generator.
            _log_turn_timing(_turn_timing_record(
                run_info, session_key, _model_ref(rec),
                text_seen=text_seen, failed=failed,
                thinking=_thinking_for_speed((rec or {}).get("speed"))))
            if draft_doc is not None:
                try:
                    update = draft_mode.post_turn_payload(draft_doc)
                    if update:
                        # NOTE: a yield in finally is deliberate (matches _DONE_SSE
                        # below). On client disconnect (GeneratorExit) the frame is
                        # silently discarded — the browser refetches the doc on
                        # session reopen, so no state is lost.
                        yield bridge._sse(update)
                except Exception:  # noqa: BLE001 - never break the turn close
                    pass
            if title_task is not None:
                try:
                    ai = await asyncio.wait_for(title_task, timeout=12)
                    if ai:
                        sessions_store.update(rec["id"], name=ai)
                except Exception:  # noqa: BLE001 - keep the first-chars fallback
                    title_task.cancel()
            # Touch the session's `updated` stamp so the semantic-search reindex
            # re-embeds this turn. It's otherwise bumped only on metadata edits
            # (title/rename/archive), so every message after a chat is first
            # titled would be invisible to search — the incremental reindex keys
            # its skip check on exactly this stamp.
            if rec:
                try:
                    sessions_store.update(rec["id"])
                except Exception:  # noqa: BLE001 - never break the turn close
                    pass
            # Auto-extract memories (gated inside: toggle pref + cooldown).
            _spawn(maybe_auto_extract(session_key))
            yield _DONE_SSE

    # Detach the turn from any single reader: a background recorder drains the
    # gateway relay into event_store and owns the turn boundary (begin/end_turn),
    # so a dropped POST reader (refresh / thread-switch / tab close) can never
    # stop or lose the run. The POST response then just TAILS event_store —
    # identical to GET /api/chat/stream — so the original POST, a thread-switch
    # return, and a post-reload resume are all the same replay-then-subscribe.
    cursor = event_store.latest_id(session_key)  # replay only THIS turn, not prior
    queue = event_store.subscribe(session_key)   # subscribe before start: no gap
    _start_turn_recorder(session_key, _drive_turn)

    async def _post_tail():
        """Replay this turn's events (everything after `cursor`), then live-tail
        new events until [DONE]. Dedupes the replay/live overlap by seq, emits a
        keepalive on idle, and always unsubscribes. Mirrors resume_route's tail
        but terminates at [DONE] (the POST closes when the turn finishes)."""
        replayed_max = -1
        try:
            for eid, payload in event_store.since(session_key, cursor):
                yield f"id: {eid}\n{payload}"
                try:
                    replayed_max = max(replayed_max, int(eid))
                except (TypeError, ValueError):
                    pass
                if _is_done_frame(payload):
                    return
            while True:
                try:
                    eid, payload = await asyncio.wait_for(
                        queue.get(), timeout=_TURN_KEEPALIVE_S)
                except asyncio.TimeoutError:
                    yield ": keepalive\n\n"  # keep proxies from killing the idle conn
                    continue
                try:
                    seq = int(eid)
                except (TypeError, ValueError):
                    seq = None
                if seq is not None and seq <= replayed_max:
                    continue  # already sent during backlog replay
                if seq is not None:
                    replayed_max = seq
                yield f"id: {eid}\n{payload}"
                if _is_done_frame(payload):
                    return
        finally:
            event_store.unsubscribe(session_key, queue)

    return StreamingResponse(_post_tail(), media_type="text/event-stream")


# --- Session persistence: metadata here, message content from the brain ------

@app.get("/api/sessions")
async def sessions():
    return sessions_store.list_sessions()


@app.post("/api/session")
async def create_session(name: str = Form(default=""), model: str = Form(default=""),
                         endpoint_url: str = Form(default=""),
                         endpoint_id: str = Form(default=""),
                         speed: str = Form(default="")):
    return sessions_store.create(name=name or None, model=model or None,
                                 endpoint_url=endpoint_url or None,
                                 endpoint_id=endpoint_id or None,
                                 speed=speed or None)


def _build_preamble(prefix: list[dict]) -> str:
    """Compact serialization of the branched-from transcript prefix, used as
    the first-send context preamble. Role + text, one per line."""
    lines = []
    for m in prefix:
        role = (m.get("role") or "user").strip()
        text = (m.get("text") or m.get("content") or "").strip()
        if not text:
            continue
        who = "Frank" if role == "user" else "Gary"
        lines.append(f"{who}: {text}")
    body = "\n".join(lines)
    return (
        "For context, this conversation was branched from an earlier thread. "
        "Here is what was said before, verbatim:\n\n"
        f"{body}\n\n"
        "Continue from here."
    )


@app.post("/api/session/branch")
async def branch_session(payload: dict = Body(default=None)):
    """Create a new session and stash a client-provided transcript prefix as
    pending context. The frontend already knows the prefix (it rendered those
    bubbles); we trust its slice verbatim and echo it back — no
    bridge.fetch_history call, no server-side re-slicing (see
    .superpowers/sdd/task-3-brief.md). The next composer submit into this
    session consumes the pending context and prepends a preamble."""
    payload = payload or {}
    source_session_id = str(payload.get("source_session_id") or "").strip()
    prefix = payload.get("prefix") or []
    if not source_session_id:
        return JSONResponse(status_code=400, content={"error": "source_session_id required"})
    if not prefix:
        return JSONResponse(status_code=400, content={"error": "prefix must not be empty"})

    src = sessions_store.get(source_session_id)
    if src is None:
        return JSONResponse(status_code=404, content={"error": "source session not found"})

    name = (payload.get("name") or "").strip() or f"↳ {src.get('name') or 'chat'}"
    new_sess = sessions_store.create(
        name=name,
        model=payload.get("model") or src.get("model"),
        endpoint_url=src.get("endpoint_url"),
        endpoint_id=src.get("endpoint_id"),
        speed=payload.get("speed") or src.get("speed"),
    )
    preamble = _build_preamble(prefix)
    branch_context.write(new_sess["id"], source_session_id, prefix, preamble)
    return {"session_id": new_sess["id"], "session_key": new_sess["sessionKey"], "prefix": prefix}


@app.get("/api/history/{session_id}")
async def history(session_id: str, limit: int = 200, cursor: str | None = None):
    """The session's saved transcript, read live from the brain. Paginated:
    the no-cursor call returns the newest `limit` messages; the frontend lazy-
    loads older windows by passing the returned `nextCursor` (see bridge
    fetch_history_page). Older pages are fetched on scroll-to-top, so no single
    response has to carry the whole transcript."""
    sess = sessions_store.get(session_id)
    if not sess:
        return {"history": [], "model": None, "hasMore": False, "nextCursor": None}
    if cursor:
        # Older pages: only the gateway's HTTP history endpoint supports
        # older-than-cursor paging.
        data = await bridge.fetch_history_page(sess["sessionKey"], limit=limit,
                                               cursor=cursor)
    else:
        # First page: the gateway's HTTP history endpoint truncates some sessions
        # — a claude-cli turn can come back as a single message, dropping Gary's
        # replies on reload. The WS chat.history returns the full transcript
        # reliably for every provider, so read the newest window from it.
        # (Tail-only: a transcript longer than the window won't lazy-load older
        # pages — an acceptable trade vs. silently losing the assistant's replies.)
        # Floor at 1000 so a reload gets the full window (not a truncated tail);
        # cap at 1000 because chat.history returns EMPTY for limits above ~1000
        # (a blank thread), so limit>1000 must not pass through.
        mapped = await bridge.fetch_history(sess["sessionKey"],
                                            limit=min(max(limit, 1000), 1000))
        data = {"history": mapped.get("history", []),
                "model": mapped.get("model"),
                "hasMore": False, "nextCursor": None}
    # use_web turns store the augmented brain message (search block + the
    # user's text) in the transcript; show only what the user typed.
    for m in data.get("history", []):
        if m.get("role") == "user":
            content = websearch.strip_context_block(m.get("content"))
            content = terminals.strip_capability_note(content)
            # A followup seed is machinery, not something Frank typed — show
            # the compact ⚙️ card line instead (frontend styles it).
            card = followup.history_card(content)
            m["content"] = card if card is not None else content
    # The terminal-attach flow records the prompt twice (the two messages differ
    # only in the stripped terminal-control note), so after stripping they're
    # identical — collapse a user message that duplicates the one just before it.
    deduped = []
    for m in data.get("history", []):
        if (m.get("role") == "user" and deduped
                and deduped[-1].get("role") == "user"
                and deduped[-1].get("content") == m.get("content")):
            continue
        deduped.append(m)
    data["history"] = deduped
    # Rehydrate image attachments persisted at send time (transcript is text-only).
    _apply_msg_attachments(session_id, deduped)
    # Prefer the record's chosen model label; fall back to whatever the brain used.
    data["model"] = sess.get("model") or data.get("model")
    return data


@app.patch("/api/session/{session_id}")
async def patch_session(session_id: str, name: str = Form(default=None),
                        model: str = Form(default=None), folder: str = Form(default=None),
                        endpoint_url: str = Form(default=None),
                        endpoint_id: str = Form(default=None),
                        speed: str = Form(default=None)):
    if speed is not None and speed not in ("fast", "normal", "deep"):
        speed = None   # invalid value → ignored, like other bad fields
    fields = {k: v for k, v in {
        "name": name, "model": model, "folder": folder,
        "endpoint_url": endpoint_url, "endpoint_id": endpoint_id,
        "speed": speed,
    }.items() if v is not None}
    return sessions_store.update(session_id, **fields) or JSONResponse(
        status_code=404, content={"detail": "no such session"})


async def _delete_gateway_session(session_key: str) -> None:
    """Best-effort gateway-side delete (transcript included) so removing a
    chat here doesn't leave its thread accumulating in the brain's session
    store — real weight on this 8GB box. Verified: sessions.delete
    {key, deleteTranscript} (deleteTranscript defaults true anyway)."""
    try:
        await bridge.gateway_call("sessions.delete",
                                  {"key": session_key, "deleteTranscript": True})
    except Exception:  # noqa: BLE001 - local delete already succeeded
        pass


@app.delete("/api/session/{session_id}")
async def delete_session(session_id: str):
    rec = sessions_store.get(session_id)
    ok = sessions_store.delete(session_id)
    if ok and rec and rec.get("sessionKey"):
        event_store.drop_session(rec["sessionKey"])  # free the in-memory event log
        _spawn(_delete_gateway_session(rec["sessionKey"]))
    return {"ok": ok}


@app.post("/api/session/{session_id}/important")
async def set_important(session_id: str, important: str = Form(default="true")):
    val = str(important).lower() not in ("false", "0", "")
    sessions_store.update(session_id, important=val)
    return {"ok": True, "important": val}


@app.post("/api/session/{session_id}/archive")
async def archive_session(session_id: str):
    sessions_store.update(session_id, archived=True)
    return {"ok": True}


@app.post("/api/session/{session_id}/unarchive")
@app.post("/api/session/{session_id}/restore")
async def unarchive_session(session_id: str):
    sessions_store.update(session_id, archived=False)
    return {"ok": True}


@app.get("/api/sessions/{session_id}/usage")
async def session_usage(session_id: str):
    """Per-session token usage + context-window weight for the footer widget.
    Thin relay over bridge.fetch_session_usage (which talks to the gateway's
    sessions.usage RPC and projects the result to the wire contract). Always
    200: on any gateway error / unknown session the body is {ok: false, reason}
    and the widget hides itself — never 500 the page."""
    return await bridge.fetch_session_usage(session_id)


# NOTE: the real resume/tail endpoints live in resume_route.py
# (/api/chat/events/resume, /api/chat/turn, /api/chat/stream). The old
# /api/chat/resume/{id} and /api/chat/stream_status/{id} stubs were dead — their
# only caller was the legacy js/sessions.js UI, which index.html no longer loads
# (the redesign app.js fully replaced it). Both stubs are removed; live run state
# is served by event_store.current_turn() via /api/chat/turn and the per-session
# working state by /api/chat/active_sessions.


@app.post("/api/chat/stop/{session_id}")
async def stop_chat(session_id: str):
    """The Stop button's server half: chat.abort the active gateway run.
    Verified shape: chat.abort {sessionKey, runId?} -> {runIds}; omitting
    runId aborts every run on the key. NOTE: on an explicit Stop the browser
    kills its fetch FIRST, which tears down gen() and pops _ACTIVE_RUNS —
    so the key-wide abort is the EXPECTED path here; the runId narrowing is
    opportunistic (e.g. stop called from another tab)."""
    session_key = sessions_store.session_key_for(session_id)
    params = {"sessionKey": session_key}
    run_id = (_ACTIVE_RUNS.get(session_key) or {}).get("runId")
    if run_id:
        params["runId"] = run_id
    try:
        payload = await bridge.gateway_call("chat.abort", params, timeout=10)
        return {"ok": True, "runIds": payload.get("runIds") or []}
    except Exception as exc:  # noqa: BLE001
        return JSONResponse(status_code=502,
                            content={"ok": False, "error": f"{exc!r}"})


@app.get("/api/models")
async def models():
    # The SPA's models.js reads `data.items` (NOT a bare array). We serve the
    # REAL gateway catalog (models.list + models.authStatus), grouped one
    # endpoint per provider (Codex / Claude), with `offline` reflecting each
    # provider's auth status. category must be "api" (models.js buckets anything
    # non-"local" there). url is the WS endpoint, echoed into the session, never
    # routed — every turn goes through the bridge regardless of picked model.
    # If the gateway is unreachable, fall back to a single honest placeholder so
    # the picker still renders (rather than going blank → "no session").
    try:
        return await bridge.fetch_models()
    except Exception:  # noqa: BLE001
        return {"items": [
            {"endpoint_id": "openclaw", "endpoint_name": "OpenClaw",
             "url": config.gateway_ws_url(), "category": "api",
             "models": ["openclaw"], "models_display": ["OpenClaw"],
             "models_extra": [], "models_extra_display": [], "offline": True},
        ]}


@app.get("/api/default-chat")
async def default_chat():
    # endpoint_url is REQUIRED: chat.js auto-creates the session only when both
    # endpoint_url and model are truthy. It's stored + echoed back to /api/session
    # but never used to route — every turn goes through the bridge regardless.
    # A user-set preference (POST /api/default-chat → .data/settings.json) wins;
    # otherwise we land on the primary agent's configured model from openclaw.json.
    pref = (websearch.load_settings().get("default_chat_model") or {})
    if pref.get("model"):
        provider = pref.get("endpoint_id") or config.default_model()[0]
        model = pref["model"]
    else:
        provider, model = config.default_model()
    return {"endpoint_id": provider, "endpoint_url": config.gateway_ws_url(),
            "model": model}


@app.post("/api/default-chat")
async def set_default_chat(payload: dict = Body(default=None)):
    """Persist the user's preferred default model for NEW web chats. Stored in
    .data/settings.json (workspace-owned) — does NOT touch the gateway's
    openclaw.json. Per-chat overrides still go through PATCH /api/session."""
    payload = payload or {}
    model = (payload.get("model") or "").strip()
    if not model:
        return JSONResponse(status_code=400, content={"error": "model required"})
    pref = {"model": model, "endpoint_id": (payload.get("endpoint_id") or "").strip()}
    websearch.save_settings({"default_chat_model": pref})
    return {"ok": True, **pref}


# Auth stubs: single-user/no-auth deployment behind Tailscale. Return a logged-in
# admin so the SPA never redirects to /login. Privileges double as feature flags:
# False hides that feature's chrome (init.js/app.js read these), so anything
# with no backend here is flipped off rather than advertised as working.
@app.get("/api/auth/status")
async def auth_status():
    return {
        "authenticated": True, "is_admin": True,
        "username": config.workspace_user(),
        "privileges": {
            "can_use_agent": True, "can_use_bash": True, "can_use_documents": True,
            "can_use_research": True,
            # No image-gen backend (hides #tool-image-btn; Gallery is CSS-hidden).
            "can_generate_images": False,
        },
    }


@app.get("/api/auth/features")
async def auth_features():
    return {"auth_required": bool(config.auth_token()), "features": {}}


@app.get("/api/auth/settings")
async def auth_settings():
    """Settings the SPA reads (search provider/result count etc.). Persisted
    in .data/settings.json; search execution lives in websearch.py."""
    return websearch.load_settings()


@app.post("/api/auth/settings")
async def save_auth_settings(payload: dict = Body(default=None)):
    return websearch.save_settings(payload or {})


@app.post("/api/search/test")
async def search_test(payload: dict = Body(default=None)):
    """One-shot probe of the configured web-search provider so Settings → Search
    can show a live OK/error. Runs a real query via websearch.py."""
    query = ((payload or {}).get("query") or "OpenClaw connectivity test").strip()
    try:
        results = await websearch.search(query, count=3)
        return {"ok": True, "count": len(results or []),
                "provider": websearch.load_settings().get("search_provider")}
    except Exception as exc:  # noqa: BLE001 — surface any provider/network error
        return {"ok": False, "error": f"{exc!r}"}


# --- Semantic search over all chat content -----------------------------------
# Real embedding-based search (Voyage voyage-3.5-lite) over every session's
# message content. Explicit routes — registered before the /api/{path} catch-all
# below and preferred by FastAPI, so they aren't shadowed by it.
@app.get("/api/search")
async def search_chats(q: str = "", limit: int = 20):
    return await chat_search.search(q, limit)


@app.post("/api/search/reindex")
async def search_reindex(force: bool = False):
    """Kick a background reindex and return immediately (a full run makes many
    gateway + Voyage calls, so it must not block the request)."""
    _spawn(chat_search.reindex(force=force))
    return {"status": "started"}


# --- Catch-all for Odysseus feature tabs v1 doesn't implement yet ------------
# cookbook, prefs, tts… each
# polls its own backend. Returning [] is universally safe: Odysseus's consumers
# all do either `data.forEach(...)` (works on []) or `data.key || []` (→ []), so
# this quiets the 404 flood without breaking any module. Registered AFTER every
# real route, so health/items/models/chat/auth/sessions still win.
@app.get("/api/{path:path}")
async def _unimplemented_api(path: str):
    return []


@app.get("/sw.js")
async def service_worker():
    """Serve the service worker from the ORIGIN ROOT.

    Registered at /static/sw.js the SW's max scope is /static/ and it can
    never control the SPA at / — the whole offline story was inert (2026-06-12
    mobile review, P0). Served at /sw.js its default scope is the origin root.
    no-cache so browsers revalidate the worker itself on each check.
    """
    sw = config.FRONTEND_DIR / "sw.js"
    if not sw.exists():
        return JSONResponse(status_code=404, content={"error": "sw.js not built"})
    return FileResponse(str(sw), media_type="application/javascript",
                        headers={"Cache-Control": "no-cache"})


# --- Serve the reused Odysseus SPA ------------------------------------------
# Mounted last so /api/* routes win. The SPA lives in frontend/ (copied from
# Odysseus static/). index.html is the entry; everything else is static assets.

def _spa_html(filename: str):
    """Serve an SPA HTML entrypoint.

    When config.BASE_PATH is set (app hosted under a stripping subpath proxy),
    rewrite browser-facing URLs so absolute /static and /api references resolve
    under the prefix: markup asset refs are rewritten, an import map remaps
    absolute dynamic imports, and a tiny network shim prefixes fetch/EventSource/
    WebSocket calls. Backend routes are untouched (the proxy strips the prefix).
    With no BASE_PATH the raw file is served byte-for-byte (default install)."""
    path = config.FRONTEND_DIR / filename
    if not path.exists():
        return JSONResponse(status_code=404, content={"error": f"{filename} not built"})
    base = config.BASE_PATH
    if not base:
        return FileResponse(str(path))
    html = path.read_text(encoding="utf-8")
    html = html.replace('="/static/', f'="{base}/static/').replace('="/api/', f'="{base}/api/')
    b = json.dumps(base)
    inject = (
        '<script type="importmap">{"imports":{"/static/":"' + base + '/static/"}}</script>'
        '<script>(function(){var B=' + b + ';'
        'window.__WS_BASE__=B;'  # asset constants (e.g. AVATAR) prefix with this
        'function fix(u){try{'
        'if(u&&typeof u==="object"&&typeof Request!=="undefined"&&u instanceof Request){return new Request(fix(u.url),u);}'
        'u=String(u);'
        'if(/^[a-z]+:\\/\\//i.test(u)){var x=new URL(u);'
        'if(x.host===location.host&&x.pathname.slice(0,B.length+1)!==B+"/"){x.pathname=B+x.pathname;}return x.toString();}'
        'if(u.slice(0,2)==="//")return u;'
        'if(u.charAt(0)==="/"&&u.slice(0,B.length+1)!==B+"/")return B+u;'
        'return u;}catch(e){return u;}}'
        'var f=window.fetch;if(f)window.fetch=function(i,n){return f(fix(i),n);};'
        'var E=window.EventSource;if(E){var NE=function(u,o){return new E(fix(u),o);};'
        'NE.prototype=E.prototype;try{NE.CONNECTING=E.CONNECTING;NE.OPEN=E.OPEN;NE.CLOSED=E.CLOSED;}catch(e){}window.EventSource=NE;}'
        'var W=window.WebSocket;if(W){var NW=function(u,p){return p===undefined?new W(fix(u)):new W(fix(u),p);};'
        'NW.prototype=W.prototype;try{NW.CONNECTING=W.CONNECTING;NW.OPEN=W.OPEN;NW.CLOSING=W.CLOSING;NW.CLOSED=W.CLOSED;}catch(e){}window.WebSocket=NW;}'
        '})();</script>'
    )
    html = html.replace("<head>", "<head>" + inject, 1)
    return HTMLResponse(html)


def _manifest_response(filename: str = "manifest.json"):
    """Serve the PWA manifest, base-path-corrected.

    Like _spa_html: when BASE_PATH is set, absolute /static icon srcs and the
    root start_url/scope must resolve under the prefix, or an installed PWA
    pulls icons from the origin root (the wrong tenant on a shared funnel) and
    launches at /. Served byte-for-byte when no BASE_PATH (default install)."""
    path = config.FRONTEND_DIR / filename
    if not path.exists():
        return JSONResponse(status_code=404, content={"error": f"{filename} not built"})
    base = config.BASE_PATH
    if not base:
        return FileResponse(str(path), media_type="application/manifest+json")
    data = json.loads(path.read_text(encoding="utf-8"))
    for icon in data.get("icons", []):
        src = icon.get("src", "")
        if isinstance(src, str) and src.startswith("/") and not src.startswith(base + "/"):
            icon["src"] = base + src
    for k in ("start_url", "scope"):
        v = data.get(k)
        if isinstance(v, str) and v.startswith("/") and not v.startswith(base + "/"):
            data[k] = base + v
    return JSONResponse(data, media_type="application/manifest+json")


if config.FRONTEND_DIR.exists():
    # Registered BEFORE the /static mount so this explicit route wins over the
    # StaticFiles handler (which would otherwise serve the raw manifest).
    @app.get("/static/manifest.json")
    async def manifest_json():
        return _manifest_response("manifest.json")

    app.mount("/static", StaticFiles(directory=str(config.FRONTEND_DIR)), name="static")

    @app.get("/")
    async def index():
        return _spa_html("index.html")

    @app.get("/classic")
    async def index_classic():
        return _spa_html("index-classic.html")
else:
    @app.get("/")
    async def index_missing():
        return JSONResponse(
            status_code=500,
            content={"error": f"frontend not found at {config.FRONTEND_DIR}. "
                              "Run scripts/sync-frontend.sh to copy Odysseus static/."},
        )
