"""Attachment subsystem for /api/chat_stream.

Extracted verbatim from app.py (Task 19): image attachments (→ gateway vision
blocks, HEIC→JPEG conversion, size caps), non-image text extraction (CSV/XLSX/
DOCX/PPTX/PDF/plain text → inlined message body), persistence of sent image refs
so /api/history can rehydrate them on reload, and terminal-drop images.

No behavior change — the chat route imports these and keeps the to_thread call
sites (the blocking ffmpeg/office/PDF work is dispatched off the event loop by
the caller in app.py, not here).
"""
from __future__ import annotations

import base64
import contextlib
import json
import logging
import mimetypes
import subprocess
import tempfile
from pathlib import Path

from . import terminals
from .fsutil import file_lock
from .uploads import ATTACH_DIR

_log = logging.getLogger(__name__)


# Gateway rejects chat.send payloads over ~10 MB base64. Cap individual images
# at 4 MB raw (~5.3 MB base64) to stay safely under that limit with multiple
# attachments. HEIC/HEIF files are converted to JPEG first (models don't support
# them natively and iPhone photos are often 7+ MB).
_ATTACH_MAX_BYTES = 4 * 1024 * 1024  # 4 MB


def _heic_to_jpeg(src: Path) -> bytes | None:
    """Convert a HEIC/HEIF file to JPEG via ffmpeg. Returns JPEG bytes or None."""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
            tmp_path = tmp.name
        result = subprocess.run(
            ["ffmpeg", "-y", "-i", str(src), "-update", "1",
             "-vf", "scale='min(1920,iw)':-2", "-q:v", "5", tmp_path],
            capture_output=True, timeout=30,
        )
        if result.returncode == 0:
            return Path(tmp_path).read_bytes()
    except Exception:  # noqa: BLE001
        pass
    finally:
        if tmp_path:
            with contextlib.suppress(OSError):
                Path(tmp_path).unlink()
    return None


def _resolve_attachments(raw: str) -> list[dict]:
    """Turn the composer's posted `attachments` (a JSON array of upload ids, each
    a filename under ATTACH_DIR) into the chat.send attachment shape the gateway
    accepts: {type, mimeType, fileName, content(base64)}. Only image/* files are
    forwarded — the gateway sniffs the bytes and drops non-images, and gpt-5.5
    takes them as inline vision blocks (large ones it offloads to media refs the
    agent resolves). Silently skips anything missing/unreadable so a bad id never
    breaks the turn."""
    if not raw:
        return []
    try:
        ids = json.loads(raw)
    except Exception:  # noqa: BLE001 - malformed field → no attachments
        return []
    out: list[dict] = []
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")  # path-traversal guard
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        if not mime.startswith("image/"):
            continue  # gateway accepts only image/* attachments
        try:
            data = path.read_bytes()
        except Exception:  # noqa: BLE001 - unreadable file → skip it
            continue
        # Convert HEIC/HEIF → JPEG (models don't support these formats).
        if mime in ("image/heic", "image/heif"):
            converted = _heic_to_jpeg(path)
            if converted:
                data, mime = converted, "image/jpeg"
            else:
                _log.warning("HEIC conversion failed for %s; skipping attachment", path.name)
                continue
        # Skip images that are still too large for the gateway after conversion.
        if len(data) > _ATTACH_MAX_BYTES:
            _log.warning("Attachment %s is %d bytes (>4 MB); skipping", path.name, len(data))
            continue
        out.append({
            "type": "image",
            "mimeType": mime,
            "fileName": path.stem + (".jpg" if mime == "image/jpeg" else path.suffix),
            "content": base64.b64encode(data).decode("ascii"),
        })
    return out


# Non-image files: extract text and inline into the message body. Total across
# all files per turn is capped so a huge upload can't blow the context.
_TEXT_ATTACH_TOTAL_MAX = 200 * 1024
_TEXT_ATTACH_PER_FILE_MAX = 100 * 1024
_TEXT_RAW_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log",
                  ".xml", ".yaml", ".yml", ".ini", ".conf", ".toml", ".env",
                  ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".htm", ".css",
                  ".scss", ".sh", ".bash", ".zsh", ".fish", ".rb", ".go", ".rs",
                  ".java", ".kt", ".swift", ".c", ".h", ".cpp", ".hpp", ".cc",
                  ".cs", ".php", ".sql", ".rst", ".org", ".vue", ".svelte",
                  ".lua", ".pl", ".r", ".dockerfile", ".makefile", ".gradle",
                  ".diff", ".patch"}


def _extract_file_text(path: Path, mime: str) -> str | None:
    """Best-effort text extraction. Returns None for formats we can't read (e.g.
    .numbers without libreoffice); the caller silently skips those."""
    ext = path.suffix.lower()
    try:
        if ext in _TEXT_RAW_EXTS or mime.startswith("text/") or mime == "application/json":
            return path.read_text(errors="replace")
        if ext == ".xlsx" or mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ):
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
            wb.close()
            return "\n".join(parts)
        if ext == ".pdf" or mime == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            out: list[str] = []
            for i, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    out.append(f"## Page {i}\n{text}")
            return "\n\n".join(out)
        if ext == ".docx" or mime == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            from docx import Document
            doc = Document(str(path))
            parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
            for tbl in doc.tables:
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
            return "\n".join(parts)
        if ext == ".pptx" or mime == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            from pptx import Presentation
            pres = Presentation(str(path))
            parts: list[str] = []
            for i, slide in enumerate(pres.slides, start=1):
                parts.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(r.text for r in para.runs).strip()
                            if line:
                                parts.append(line)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            if any(cells):
                                parts.append("\t".join(cells))
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Speaker notes] {notes}")
            return "\n".join(parts)
    except Exception as e:  # noqa: BLE001
        _log.warning("Text extraction failed for %s (%s): %s", path.name, ext, e)
        return None
    return None


def _extract_text_attachments(raw: str) -> tuple[list[dict], list[str]]:
    """Return (files, skipped). `files` is [{name, text}] for uploads we could
    read; `skipped` is names of non-image uploads whose content we couldn't
    extract (unsupported format, empty output, or budget exceeded). Images are
    handled by `_resolve_attachments` and are neither returned nor skipped
    here."""
    if not raw:
        return [], []
    try:
        ids = json.loads(raw)
    except Exception:  # noqa: BLE001
        return [], []
    out: list[dict] = []
    skipped: list[str] = []
    total = 0
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or ""
        if mime.startswith("image/"):
            continue  # image path handles these
        text = _extract_file_text(path, mime)
        if not text or not text.strip():
            skipped.append(path.name)
            continue
        text = text.strip()
        if len(text) > _TEXT_ATTACH_PER_FILE_MAX:
            text = text[:_TEXT_ATTACH_PER_FILE_MAX] + "\n... [truncated]"
        if total + len(text) > _TEXT_ATTACH_TOTAL_MAX:
            _log.warning("Text-attach total exceeded %d bytes; dropping %s",
                         _TEXT_ATTACH_TOTAL_MAX, path.name)
            skipped.append(path.name)
            continue
        total += len(text)
        out.append({"name": path.name, "text": text})
    return out, skipped


def _prepend_text_attachments(message: str, files: list[dict],
                              skipped: list[str] | None = None) -> str:
    """Inline extracted file text into the user message so it reaches the model
    regardless of vision/file-type support. The gateway only forwards image
    attachments; everything else has to arrive as text."""
    skipped = skipped or []
    if not files and not skipped:
        return message
    parts: list[str] = []
    if files:
        parts.append("The user attached the following file(s):\n")
        for f in files:
            parts.append(f"── {f['name']} ──")
            parts.append(f["text"])
            parts.append("")
        parts.append("──\n")
    if skipped:
        parts.append(
            "[System note: couldn't extract text from these attachments — "
            "unsupported format or empty content: "
            + ", ".join(skipped) + ". Tell the user which files were dropped "
            "so they can convert or resend.]\n"
        )
    parts.append(message or "")
    return "\n".join(parts)


_CHAT_ATTACH_DIR = ATTACH_DIR.parent / ".chat-attachments"


def _attach_log_path(session_id: str) -> Path | None:
    safe = "".join(c for c in (session_id or "") if c.isalnum() or c in "-_.")
    return (_CHAT_ATTACH_DIR / f"{safe}.json") if safe else None


def _persist_msg_attachments(session_id: str, message: str, attachments_raw: str) -> None:
    """Record image attachments for a sent message so /api/history can rehydrate
    them on reload — the gateway transcript keeps only text, so a sent image
    otherwise vanishes on refresh. Appends one record per send-with-images:
    {text, att:[{id,url}]}. Best-effort; never raises into the turn."""
    if not session_id or not attachments_raw:
        return
    try:
        ids = json.loads(attachments_raw)
    except Exception:  # noqa: BLE001 - malformed field → nothing to persist
        return
    att: list[dict] = []
    for fid in ids if isinstance(ids, list) else []:
        if not isinstance(fid, str):
            continue
        safe = "".join(c for c in fid if c.isalnum() or c in "-_.")
        path = ATTACH_DIR / safe
        if not path.exists() or not path.is_file():
            continue
        mime = mimetypes.guess_type(str(path))[0] or ""
        if not mime.startswith("image/"):
            continue
        att.append({"id": safe, "url": f"/api/upload/{safe}"})
    if not att:
        return
    p = _attach_log_path(session_id)
    if not p:
        return
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        with file_lock(p):
            log = []
            if p.exists():
                try:
                    log = json.loads(p.read_text() or "[]")
                except Exception:  # noqa: BLE001
                    log = []
            log.append({"text": (message or "").strip(), "att": att})
            p.write_text(json.dumps(log))
    except Exception:  # noqa: BLE001 - persistence is best-effort
        pass


def _apply_msg_attachments(session_id: str, history: list[dict]) -> None:
    """Assign persisted image attachments to user messages in `history`, matching
    on (trimmed) text in send order. Mutates `history` in place; best-effort."""
    p = _attach_log_path(session_id)
    if not p or not p.exists():
        return
    try:
        log = json.loads(p.read_text() or "[]")
    except Exception:  # noqa: BLE001
        return
    if not isinstance(log, list) or not log:
        return
    used = [False] * len(log)
    for m in history:
        if m.get("role") != "user":
            continue
        text = (m.get("content") or "").strip()
        for i, rec in enumerate(log):
            if used[i]:
                continue
            if (rec.get("text") or "").strip() == text:
                m["attachments"] = rec.get("att") or []
                used[i] = True
                break


def _terminal_attachments(terminal_key: str) -> list[dict]:
    """Pending images the user dropped into this chat's terminal → chat.send
    image blocks, then mark them consumed so each rides exactly one turn. The
    token→path mapping itself persists (terminals registry) for later resolves.
    Mirrors _resolve_attachments' block shape; image/* only; bad files skipped."""
    out: list[dict] = []
    consumed: list[str] = []
    for it in terminals.list_attachments(terminal_key, pending_only=True):
        path = Path(it.get("path", ""))
        if not path.is_file():
            continue
        mime = it.get("mime") or mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        if not mime.startswith("image/"):
            continue
        try:
            data = path.read_bytes()
        except Exception:  # noqa: BLE001 - unreadable → skip, never break the turn
            continue
        out.append({
            "type": "image",
            "mimeType": mime,
            "fileName": path.name,
            "content": base64.b64encode(data).decode("ascii"),
        })
        consumed.append(it["token"])
    if consumed:
        terminals.mark_consumed(terminal_key, consumed)
    return out
